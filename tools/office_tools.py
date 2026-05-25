import os
import sys

def docx_to_pdf(docx_path: str, pdf_path: str = None) -> str:
    from tools.doc_tools import word_to_pdf
    return word_to_pdf(docx_path, pdf_path)


def pdf_to_docx(pdf_path: str, docx_path: str = None) -> str:
    from tools.doc_tools import pdf_to_word
    return pdf_to_word(pdf_path, docx_path)


def xlsx_to_pdf(xlsx_path: str, pdf_path: str = None) -> str:
    """Enhanced Excel to PDF conversion using pandas and Qt's rendering engine."""
    try:
        import pandas as pd
        from PySide6.QtGui import QTextDocument, QPdfWriter, QPageLayout, QPageSize
        from PySide6.QtCore import QMarginsF

        if pdf_path is None:
            pdf_path = os.path.splitext(xlsx_path)[0] + ".pdf"

        xlsx = pd.ExcelFile(xlsx_path)
        all_html = []
        for sheet_name in xlsx.sheet_names:
            df = pd.read_excel(xlsx, sheet_name=sheet_name)
            # Fill NaN with empty string
            df = df.fillna("")
            html = df.to_html(index=False, border=1)
            all_html.append(f"<h2>Sheet: {sheet_name}</h2>{html}")
        
        combined_html = f"""
        <html>
        <head>
        <style>
            body {{ font-family: sans-serif; font-size: 9pt; margin: 20px; }}
            table {{ border-collapse: collapse; width: 100%; margin-bottom: 25px; table-layout: auto; }}
            th, td {{ border: 1px solid #666; padding: 6px; text-align: left; }}
            th {{ background-color: #eee; font-weight: bold; }}
            h2 {{ color: #2c3e50; border-bottom: 1px solid #ccc; padding-bottom: 5px; }}
        </style>
        </head>
        <body>{"".join(all_html)}</body>
        </html>
        """

        doc = QTextDocument()
        doc.setHtml(combined_html)

        writer = QPdfWriter(pdf_path)
        writer.setPageSize(QPageSize(QPageSize.PageSizeId.A4))
        writer.setPageOrientation(QPageLayout.Orientation.Landscape)
        writer.setPageMargins(QMarginsF(10, 10, 10, 10), QPageLayout.Unit.Millimeter)
        
        doc.print_(writer)
        return pdf_path
    except Exception as e:
        print(f"Advanced Excel to PDF failed: {e}. Using fallback...")
        import openpyxl
        from fpdf import FPDF

        if pdf_path is None:
            pdf_path = os.path.splitext(xlsx_path)[0] + ".pdf"

        wb = openpyxl.load_workbook(xlsx_path)
        ws = wb.active
        pdf = FPDF()
        pdf.add_page()
        pdf.set_font("Courier", "", 8)
        for row in ws.iter_rows(values_only=True):
            line = " | ".join(str(cell) if cell is not None else "" for cell in row)
            pdf.multi_cell(0, 5, line)
        pdf.output(pdf_path)
        return pdf_path


def pdf_to_xlsx(pdf_path: str, xlsx_path: str = None) -> str:
    """Real table extraction from PDF to Excel using pdfplumber."""
    try:
        import pdfplumber
        import pandas as pd

        if xlsx_path is None:
            xlsx_path = os.path.splitext(pdf_path)[0] + ".xlsx"

        with pdfplumber.open(pdf_path) as pdf:
            with pd.ExcelWriter(xlsx_path, engine='openpyxl') as writer:
                for i, page in enumerate(pdf.pages):
                    table = page.extract_table()
                    if table:
                        df = pd.DataFrame(table[1:], columns=table[0])
                        df.to_excel(writer, sheet_name=f'Page_{i+1}', index=False)
        return xlsx_path
    except Exception as e:
        print(f"Advanced PDF to Excel failed: {e}. Using fallback...")
        import openpyxl
        import pypdf
        if xlsx_path is None:
            xlsx_path = os.path.splitext(pdf_path)[0] + ".xlsx"
        wb = openpyxl.Workbook()
        ws = wb.active
        reader = pypdf.PdfReader(pdf_path)
        row = 1
        for page in reader.pages:
            text = page.extract_text()
            if text:
                for line in text.splitlines():
                    if line.strip():
                        ws.cell(row=row, column=1, value=line.strip())
                        row += 1
        wb.save(xlsx_path)
        return xlsx_path


def pptx_to_pdf(pptx_path: str, pdf_path: str = None) -> str:
    """Robust PPTX to PDF using text extraction and Qt rendering."""
    try:
        from pptx import Presentation
        from PySide6.QtGui import QTextDocument, QPdfWriter, QPageLayout, QPageSize
        from PySide6.QtCore import QMarginsF
        
        if pdf_path is None:
            pdf_path = os.path.splitext(pptx_path)[0] + ".pdf"

        prs = Presentation(pptx_path)
        all_html = []
        for i, slide in enumerate(prs.slides):
            slide_content = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        if para.text.strip():
                            slide_content.append(f"<p>{para.text.strip()}</p>")
            
            content = "".join(slide_content)
            all_html.append(f"<div class='slide'><h2>Slide {i+1}</h2>{content}</div>")

        combined_html = f"""
        <html>
        <head>
        <style>
            body {{ font-family: sans-serif; margin: 40px; }}
            .slide {{ page-break-after: always; border: 1px solid #ccc; padding: 20px; margin-bottom: 20px; }}
            h2 {{ color: #2980b9; }}
        </style>
        </head>
        <body>{"".join(all_html)}</body>
        </html>
        """

        doc = QTextDocument()
        doc.setHtml(combined_html)
        writer = QPdfWriter(pdf_path)
        writer.setPageSize(QPageSize(QPageSize.PageSizeId.A4))
        writer.setPageMargins(QMarginsF(15, 15, 15, 15), QPageLayout.Unit.Millimeter)
        doc.print_(writer)
        return pdf_path
    except Exception as e:
        print(f"PPTX to PDF failed: {e}. Using fallback...")
        from pptx import Presentation
        from fpdf import FPDF
        if pdf_path is None:
            pdf_path = os.path.splitext(pptx_path)[0] + ".pdf"
        prs = Presentation(pptx_path)
        pdf = FPDF()
        pdf.set_font("Courier", "", 10)
        for slide in prs.slides:
            pdf.add_page()
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        if para.text.strip():
                            pdf.multi_cell(0, 6, para.text.strip())
        pdf.output(pdf_path)
        return pdf_path


def pdf_to_pptx(pdf_path: str, pptx_path: str = None) -> str:
    """PDF to PPTX using text extraction."""
    from pptx import Presentation
    from pptx.util import Inches
    import pypdf

    if pptx_path is None:
        pptx_path = os.path.splitext(pdf_path)[0] + ".pptx"

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    reader = pypdf.PdfReader(pdf_path)
    for page in reader.pages:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        text = page.extract_text() or ""
        txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), prs.slide_width - Inches(1), prs.slide_height - Inches(1))
        tf = txBox.text_frame
        tf.word_wrap = True
        tf.text = text
    prs.save(pptx_path)
    return pptx_path


def csv_to_xlsx(csv_path: str, xlsx_path: str = None, delimiter: str = ",") -> str:
    import pandas as pd
    if xlsx_path is None:
        xlsx_path = os.path.splitext(csv_path)[0] + ".xlsx"
    df = pd.read_csv(csv_path, sep=delimiter)
    df.to_excel(xlsx_path, index=False)
    return xlsx_path


def xlsx_to_csv(xlsx_path: str, csv_path: str = None) -> str:
    import pandas as pd
    if csv_path is None:
        csv_path = os.path.splitext(xlsx_path)[0] + ".csv"
    df = pd.read_excel(xlsx_path)
    df.to_csv(csv_path, index=False)
    return csv_path


def image_to_pdf(image_paths: list, pdf_path: str = "output.pdf") -> str:
    import img2pdf
    with open(pdf_path, "wb") as f:
        f.write(img2pdf.convert(image_paths))
    return pdf_path


def text_to_pdf(text: str, pdf_path: str = "output.pdf") -> str:
    from tools.doc_tools import text_to_pdf as t2p
    return t2p(text, pdf_path)


def merge_pdfs(pdf_list: list, output_path: str = "merged.pdf") -> str:
    from tools.pdf import merge_documents
    merge_documents(pdf_list, output_path)
    return output_path


def split_pdf(pdf_path: str, output_dir: str = None) -> list:
    from tools.doc_tools import split_pdf as sp
    return sp(pdf_path, output_dir)


def pdf_to_image(pdf_path: str, output_dir: str = None, fmt: str = "png", dpi: int = 200) -> list:
    import fitz
    if output_dir is None:
        output_dir = os.path.dirname(pdf_path) or "."
    os.makedirs(output_dir, exist_ok=True)
    doc = fitz.open(pdf_path)
    base = os.path.splitext(os.path.basename(pdf_path))[0]
    paths = []
    for page_num in range(len(doc)):
        page = doc[page_num]
        pix = page.get_pixmap(dpi=dpi)
        out_path = os.path.join(output_dir, f"{base}_page_{page_num+1}.{fmt}")
        pix.save(out_path)
        paths.append(out_path)
    doc.close()
    return paths


def pdf_to_text(pdf_path: str, output_path: str = None) -> str:
    import pypdf
    if output_path is None:
        output_path = os.path.splitext(pdf_path)[0] + ".txt"
    reader = pypdf.PdfReader(pdf_path)
    text = "".join(page.extract_text() or "" for page in reader.pages)
    with open(output_path, "w", encoding="utf-8") as f:
        f.write(text)
    return output_path
